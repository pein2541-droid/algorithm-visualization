from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)       # dark navy
PRIMARY = RGBColor(0x14, 0xB8, 0xA6)        # teal-500
PRIMARY_DARK = RGBColor(0x0D, 0x94, 0x88)   # teal-600
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)     # slate-400
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)         # amber-500
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)      # slate-800
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_shape_rect(slide, left, top, width, height, fill_color=DARK_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_accent_bar(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

def section_header(slide, number, title, subtitle=""):
    """Standard section divider slide"""
    add_bg(slide)
    # Accent line
    add_accent_bar(slide, 1.5, 3.2, 0.8)
    # Section number
    add_text_box(slide, 1.5, 2.2, 2, 0.8, f'PART {number}', font_size=16, color=PRIMARY, bold=True)
    # Title
    add_text_box(slide, 1.5, 3.4, 10, 1.2, title, font_size=44, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, 1.5, 4.6, 10, 0.6, subtitle, font_size=18, color=LIGHT_GRAY)

def content_slide(slide, title_text, title_num=""):
    """Standard content slide with title bar"""
    add_bg(slide)
    # Top bar
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x1E)
    header_bar.line.fill.background()
    # Accent underline
    add_accent_bar(slide, 0, 1.1, 13.333)
    # Title
    prefix = f'0{title_num}  ' if title_num else ''
    add_text_box(slide, 1, 0.15, 11, 0.8, f'{prefix}{title_text}', font_size=28, color=WHITE, bold=True)

# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
# Decorative accent shape
dec = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
dec.fill.solid(); dec.fill.fore_color.rgb = PRIMARY; dec.line.fill.background()

add_text_box(slide, 1.5, 1.8, 10, 0.5, '毕业设计答辩', font_size=20, color=PRIMARY, bold=True)
add_text_box(slide, 1.5, 2.5, 10, 1.5, 'Python 算法动画演示系统', font_size=52, color=WHITE, bold=True)
add_text_box(slide, 1.5, 4.2, 10, 0.5, 'Algorithm Visualization Learning Platform', font_size=24, color=LIGHT_GRAY)
add_text_box(slide, 1.5, 5.5, 10, 0.5, '基于 React + Flask + PostgreSQL 的全栈算法可视化教学平台', font_size=18, color=LIGHT_GRAY)

# Bottom info
add_text_box(slide, 1.5, 6.5, 4, 0.4, '答辩人：XXX', font_size=16, color=LIGHT_GRAY)
add_text_box(slide, 8, 6.5, 4, 0.4, '指导教师：XXX', font_size=16, color=LIGHT_GRAY)

# ============================================================
# SLIDE 2 — 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.8, '目  录', font_size=40, color=WHITE, bold=True)
add_accent_bar(slide, 1.5, 1.6, 1.5)

items = [
    ('01', '需求分析', '项目背景、用户角色、功能需求、技术选型'),
    ('02', '系统设计', '系统架构、数据库设计、API 设计、前端组件设计'),
    ('03', '运行演示', '核心功能演示、页面展示、操作流程'),
    ('04', '总结', '项目成果、创新点、不足与展望'),
    ('05', '致谢', '感谢指导教师与评审专家'),
]
for i, (num, title, desc) in enumerate(items):
    y = 2.2 + i * 1.0
    add_text_box(slide, 2.0, y, 0.8, 0.5, num, font_size=28, color=PRIMARY, bold=True)
    add_text_box(slide, 3.0, y, 3, 0.5, title, font_size=24, color=WHITE, bold=True)
    add_text_box(slide, 3.0, y + 0.4, 8, 0.4, desc, font_size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 3 — Section: 需求分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 1, '需求分析', 'Requirements Analysis')

# ============================================================
# SLIDE 4 — 项目背景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '项目背景与问题分析', '1')

points = [
    ('算法学习痛点',
     '传统算法教学以静态教材和PPT为主，学生难以直观理解算法的执行过程。\n排序、搜索、树遍历等算法的动态执行逻辑仅凭文字和图片难以形成有效认知。'),
    ('缺乏交互性',
     '现有教学平台多为视频播放，缺少逐步骤的交互控制能力。\n学生无法自主控制学习节奏，不能暂停、回退到特定的执行步骤。'),
    ('理论与实践脱节',
     '伪代码与实际执行过程之间缺乏同步映射机制。\n学生在看代码时，难以对应到每一步的具体操作（如比较、交换、递归）。'),
]
for i, (title, desc) in enumerate(points):
    y = 1.6 + i * 1.8
    add_shape_rect(slide, 1.2, y, 0.08, 1.2, fill_color=PRIMARY)
    add_text_box(slide, 1.8, y, 10, 0.5, title, font_size=22, color=PRIMARY, bold=True)
    add_text_box(slide, 1.8, y + 0.5, 10, 0.9, desc, font_size=15, color=LIGHT_GRAY)

# ============================================================
# SLIDE 5 — 用户角色
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '用户角色与核心需求', '2')

roles = [
    ('学生', '👤', '浏览算法列表\n查看动画演示\n学习理论知识\n记录学习进度\n收藏感兴趣算法', BLUE),
    ('教师', '👩‍🏫', '浏览所有算法\n使用动画辅助教学\n查看学习资料\n管理个人资料', GREEN),
    ('管理员', '🛡️', '用户管理（启用/禁用/删除）\n算法管理（增删改查）\n分类与学习材料管理\n系统数据统计仪表盘', RED),
]
for i, (role, icon, desc, color) in enumerate(roles):
    x = 1.2 + i * 3.9
    card = add_shape_rect(slide, x, 1.8, 3.5, 4.8, fill_color=DARK_CARD, border_color=color)
    add_text_box(slide, x + 0.3, 2.0, 3, 0.4, icon, font_size=32, color=WHITE)
    add_text_box(slide, x + 0.3, 2.6, 3, 0.5, role, font_size=24, color=color, bold=True)
    add_text_box(slide, x + 0.3, 3.3, 3, 3.0, desc, font_size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 6 — 功能模块
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '功能模块划分', '3')

add_text_box(slide, 1.2, 1.6, 10, 0.4, '系统功能模块总览', font_size=20, color=PRIMARY, bold=True)

modules = [
    ('算法浏览与筛选', '按排序/搜索/树/图/DP 分类浏览\n算法名称、复杂度、描述卡片展示', '📊'),
    ('动画演示核心', '逐步骤动画播放/暂停\n单步前进/后退，速度调节\n伪代码高亮同步映射', '🎬'),
    ('学习资料系统', '分类学习材料\nHTML 富文本内容\n多语言代码示例', '📚'),
    ('用户系统', 'JWT 注册/登录\n学习记录追踪\n算法收藏功能', '👤'),
    ('管理后台', '用户管理\n算法 CRUD\n分类与材料管理\n数据统计仪表盘', '⚙️'),
    ('辅助功能', '暗色/亮色主题切换\n新手引导教程\n响应式布局', '🌓'),
]
for i, (title, desc, icon) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = 1.2 + col * 4
    y = 2.3 + row * 2.5
    card = add_shape_rect(slide, x, y, 3.6, 2.2, fill_color=DARK_CARD)
    add_text_box(slide, x + 0.3, y + 0.15, 3, 0.4, icon, font_size=24, color=WHITE)
    add_text_box(slide, x + 0.3, y + 0.55, 3, 0.4, title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, x + 0.3, y + 1.0, 3, 1.0, desc, font_size=13, color=LIGHT_GRAY)

# ============================================================
# SLIDE 7 — Section: 系统设计
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 2, '系统设计', 'System Design')

# ============================================================
# SLIDE 8 — 系统架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '系统架构设计', '1')

# Architecture diagram using shapes
layers = [
    ('展示层 (Frontend)', 'React 18 + TypeScript\nTailwindCSS + Three.js\nZustand 状态管理', PRIMARY),
    ('通信层 (Communication)', 'RESTful API\nAxios + JWT Token\nVite Proxy 代理', ACCENT),
    ('业务层 (Backend)', 'Python Flask\nFlask-JWT-Extended 认证\nService 服务层模式', RGBColor(0x3B, 0x82, 0xF6)),
    ('数据层 (Database)', 'PostgreSQL (Neon)\nSQLAlchemy ORM\n7 张数据表', RGBColor(0x22, 0xC5, 0x5E)),
]
for i, (name, desc, color) in enumerate(layers):
    y = 1.6 + i * 1.4
    shape = add_shape_rect(slide, 3.0, y, 7.5, 1.15, fill_color=DARK_CARD, border_color=color)
    add_text_box(slide, 3.3, y + 0.1, 2.2, 0.35, name, font_size=16, color=color, bold=True)
    add_text_box(slide, 3.3, y + 0.45, 6.8, 0.6, desc, font_size=13, color=LIGHT_GRAY)
    # Arrow down (except last)
    if i < len(layers) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.4), Inches(y + 1.15), Inches(0.35), Inches(0.25))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = LIGHT_GRAY; arrow.line.fill.background()

# ============================================================
# SLIDE 9 — 技术选型
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '技术选型说明', '2')

techs = [
    ('前端技术', [
        'React 18 — 主流前端框架，组件化开发，生态成熟',
        'TypeScript — 类型安全，减少运行时错误',
        'Vite 4 — 极速开发服务器，HMR 热更新',
        'TailwindCSS — 原子化 CSS，快速构建 UI',
        'Zustand — 轻量状态管理，persist 持久化',
        'Three.js / React-Three-Fiber — 3D 可视化渲染',
    ]),
    ('后端技术', [
        'Flask — 轻量 Python Web 框架，灵活可扩展',
        'Flask-JWT-Extended — JWT 认证，角色权限控制',
        'SQLAlchemy ORM — 数据库抽象层，迁移友好',
        'Marshmallow — 数据序列化/校验',
        'bcrypt — 密码安全哈希',
    ]),
    ('基础设施', [
        'PostgreSQL (Neon) — 云托管关系型数据库',
        'Vite Proxy — 开发环境前后端代理',
        'concurrently — 前后端并行启动',
        'Git — 版本控制',
    ]),
]
for i, (category, items) in enumerate(techs):
    x = 1.2 + i * 3.8
    add_text_box(slide, x, 1.5, 3.5, 0.5, category, font_size=20, color=PRIMARY, bold=True)
    for j, item in enumerate(items):
        add_text_box(slide, x, 2.2 + j * 0.55, 3.5, 0.5, f'• {item}', font_size=13, color=LIGHT_GRAY)

# ============================================================
# SLIDE 10 — 数据库设计
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '数据库设计', '3')

tables = [
    ('users', 'user_id, username, email, password_hash, role_type, status, created_at', '用户信息与角色'),
    ('algorithms', 'algorithm_id, name, type, description, time_complexity, space_complexity, pseudocode(JSON)', '算法元数据'),
    ('animation_data', 'animation_id, algorithm_id(FK), animation_steps(JSON), initial_data(JSON), visualization_type', '动画步骤数据'),
    ('categories', 'category_id, name, description, sort_order', '学习材料分类'),
    ('learning_materials', 'material_id, category_id(FK), algorithm_id(FK), title, content, code_example(JSON)', '学习内容'),
    ('learning_records', 'record_id, user_id(FK), algorithm_id(FK), duration_seconds, progress_status', '用户学习记录'),
    ('favorite_algorithms', 'favorite_id, user_id(FK), algorithm_id(FK), favorited_at', '用户收藏'),
]
# Table header
add_text_box(slide, 1.2, 1.5, 10, 0.4, 'E-R 关系：User → LearningRecord/Favorite ← Algorithm → AnimationData/LearningMaterial ← Category', font_size=14, color=LIGHT_GRAY)

for i, (table, fields, note) in enumerate(tables):
    y = 2.1 + i * 0.72
    add_shape_rect(slide, 1.2, y, 11, 0.6, fill_color=DARK_CARD if i % 2 == 0 else RGBColor(0x16, 0x20, 0x30))
    add_text_box(slide, 1.4, y + 0.1, 2.0, 0.4, table, font_size=14, color=PRIMARY, bold=True)
    add_text_box(slide, 3.5, y + 0.1, 6.5, 0.4, fields, font_size=11, color=WHITE)
    add_text_box(slide, 10.2, y + 0.1, 2.5, 0.4, note, font_size=11, color=LIGHT_GRAY)

# ============================================================
# SLIDE 11 — API 设计
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, 'RESTful API 设计', '4')

api_groups = [
    ('认证接口', 'POST /api/auth/register\nPOST /api/auth/login', '公开'),
    ('算法接口', 'GET /api/algorithms\nGET /api/algorithms/:id\nGET /api/algorithms/:id/animation', '公开'),
    ('学习接口', 'GET /api/learning/categories\nGET /api/learning/categories/:id/content', '公开'),
    ('用户接口', 'GET /api/user/profile', '需登录'),
    ('管理接口 (15个)', 'GET /api/admin/stats\nGET/PUT/DELETE /api/admin/users\nPOST/PUT/DELETE /api/admin/algorithms\nGET/POST/PUT/DELETE /api/admin/categories\nGET/POST/PUT/DELETE /api/admin/materials', '需管理员'),
]
for i, (group, routes, auth) in enumerate(api_groups):
    x = 1.2 if i < 3 else 6.8
    y = 1.6 + (i % 3) * 1.85
    card = add_shape_rect(slide, x, y, 5.3, 1.65, fill_color=DARK_CARD)
    add_text_box(slide, x + 0.3, y + 0.1, 3, 0.35, group, font_size=17, color=PRIMARY, bold=True)
    add_text_box(slide, x + 0.3, y + 0.5, 4.5, 1.0, routes, font_size=12, color=WHITE)
    auth_color = GREEN if auth == '公开' else (ACCENT if auth == '需登录' else RED)
    add_text_box(slide, x + 4.0, y + 0.1, 1.0, 0.35, auth, font_size=12, color=auth_color, bold=True)

# ============================================================
# SLIDE 12 — 前端组件设计
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '前端组件架构', '5')

# Component tree
tree = [
    ('App.tsx (路由根)', 0),
    ('├─ Home (首页)', 1), ('├─ Algorithms (算法列表)', 1), ('├─ AlgorithmDetail (算法详情)', 1),
    ('├─ Learning (学习资料)', 1), ('├─ Profile (个人中心)', 1), ('├─ Login / Register (认证)', 1),
    ('├─ AdminLayout (管理布局)', 1), ('│  ├─ AdminDashboard (仪表盘)', 2),
    ('│  ├─ AdminUsers (用户管理)', 2), ('│  ├─ AdminAlgorithms (算法管理)', 2),
    ('│  └─ AdminContent (内容管理)', 2),
    ('Zustand Store (状态管理)', 0),
    ('├─ useAuthStore (认证 + persist)', 1), ('├─ useThemeStore (主题 + persist)', 1),
    ('└─ useAnimationStore (动画播放)', 1),
]
for i, (text, indent) in enumerate(tree):
    x = 1.2 + indent * 1.5
    y = 1.6 + i * 0.37
    color = PRIMARY if indent == 0 else LIGHT_GRAY
    bold = indent == 0
    add_text_box(slide, x, y, 8, 0.35, text, font_size=14, color=color, bold=bold)

# Right side: Visualization components
add_text_box(slide, 8.5, 1.6, 4, 0.4, '可视化组件 (6个)', font_size=16, color=PRIMARY, bold=True)
viz = ['BubbleSortComponent', 'KnapsackComponent', 'DFSComponent',
       'BFSComponent', 'BinaryTreeComponent', 'SearchComponent']
for i, v in enumerate(viz):
    add_text_box(slide, 8.5, 2.2 + i * 0.45, 4, 0.35, f'• {v}', font_size=13, color=LIGHT_GRAY)

# ============================================================
# SLIDE 13 — Section: 运行演示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 3, '运行演示', 'Live Demonstration')

# ============================================================
# SLIDE 14 — 演示思路
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '演示流程设计', '1')

demo_steps = [
    ('1', '系统首页', '展示 Landing Page\na) Hero 区域自动切换动画预览\nb) 四大核心特色卡片\nc) 热门算法展示', PRIMARY),
    ('2', '算法浏览与筛选', '展示 /algorithms 页面\na) 按 5 种类型分类筛选\nb) 算法卡片响应式布局\nc) 时间/空间复杂度展示', BLUE),
    ('3', '核心：动画演示', '展示算法详情页（以冒泡排序为例）\na) 播放/暂停动画\nb) 单步前进/后退\nc) 速度调节（0.5x - 3x）\nd) 伪代码高亮同步', GREEN),
    ('4', '学习资料', '展示 /learning 页面\na) 分类侧边导航\nb) HTML 富文本内容\nc) 代码示例复制', PURPLE),
    ('5', '用户系统', '展示登录/注册/个人中心\na) JWT 注册登录流程\nb) 学习记录追踪\nc) 算法收藏', ACCENT),
    ('6', '管理后台', '展示 /admin 管理功能\na) 仪表盘统计\nb) 用户管理（禁用/删除）\nc) 算法/材料 CRUD', RED),
]
for i, (num, title, desc, color) in enumerate(demo_steps):
    col = i % 2
    row = i // 2
    x = 1.2 + col * 6.2
    y = 1.6 + row * 1.85
    card = add_shape_rect(slide, x, y, 5.6, 1.65, fill_color=DARK_CARD, border_color=color)
    add_text_box(slide, x + 0.3, y + 0.1, 0.5, 0.4, num, font_size=28, color=color, bold=True)
    add_text_box(slide, x + 1.0, y + 0.1, 4, 0.4, title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, x + 0.3, y + 0.55, 5, 1.0, desc, font_size=13, color=LIGHT_GRAY)

# ============================================================
# SLIDE 15 — 核心亮点截图说明
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '核心功能亮点', '2')

highlights = [
    ('逐步骤动画引擎', '后端 AnimationService 为每种算法\ntype 自动生成动画步骤 JSON',
     '排序 → 比较/交换序列\nDP → 表格填充步骤\n图 → 节点访问顺序\n树 → 遍历路径'),
    ('伪代码同步高亮', '当前执行步骤与伪代码行\n实时映射，同步高亮显示',
     '动画步骤索引 →\n伪代码行号映射 →\n语法高亮着色'),
    ('角色权限体系', '三层角色（student/teacher/admin）\nJWT + admin_required 装饰器',
     '公开接口：浏览/查询\n认证接口：profile\n管理员接口：15 个 CRUD'),
    ('暗色/亮色主题', 'TailwindCSS darkMode: class\nZustand persist 持久化',
     '全站 30+ 组件\n完整的暗色适配\n用户偏好记忆'),
]
for i, (title, desc, detail) in enumerate(highlights):
    col = i % 2
    row = i // 2
    x = 1.2 + col * 6.2
    y = 1.6 + row * 2.8
    card = add_shape_rect(slide, x, y, 5.6, 2.5, fill_color=DARK_CARD)
    add_text_box(slide, x + 0.3, y + 0.15, 5, 0.4, title, font_size=20, color=PRIMARY, bold=True)
    add_text_box(slide, x + 0.3, y + 0.65, 2.4, 1.2, desc, font_size=13, color=WHITE)
    add_text_box(slide, x + 2.9, y + 0.65, 2.5, 1.6, detail, font_size=13, color=LIGHT_GRAY)

# ============================================================
# SLIDE 16 — Section: 总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 4, '总结与展望', 'Summary & Future Work')

# ============================================================
# SLIDE 17 — 项目成果
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '项目成果与技术指标', '1')

# Stats row
stats_data = [
    ('9', '核心算法', PRIMARY),
    ('6', '可视化组件', BLUE),
    ('25', 'API 端点', GREEN),
    ('7', '数据表', PURPLE),
    ('10', '前端页面', ACCENT),
    ('3', '用户角色', RED),
]
for i, (num, label, color) in enumerate(stats_data):
    x = 1.2 + i * 1.95
    card = add_shape_rect(slide, x, 1.6, 1.75, 1.6, fill_color=DARK_CARD, border_color=color)
    add_text_box(slide, x + 0.1, 1.8, 1.5, 0.7, num, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.1, 2.5, 1.5, 0.4, label, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key achievements
add_text_box(slide, 1.2, 3.6, 10, 0.5, '主要成果', font_size=22, color=PRIMARY, bold=True)
achievements = [
    '完成了前后端分离的全栈 Web 应用开发，覆盖算法演示、学习资料、用户管理完整功能闭环',
    '设计并实现了逐步骤动画引擎，支持排序、搜索、树、图、DP 五大类算法可视化演示',
    '实现了伪代码与动画步骤的实时同步高亮机制，增强算法学习的直观性',
    '构建了基于 JWT 的三层角色权限体系（学生/教师/管理员），管理员后台支持完整的数据管理',
    '采用现代前端工程化方案（React + TypeScript + Vite + TailwindCSS），代码类型安全、构建高效',
]
for i, ach in enumerate(achievements):
    add_text_box(slide, 1.2, 4.3 + i * 0.48, 11, 0.45, f'{i+1}. {ach}', font_size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 18 — 创新点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '创新点与特色', '2')

innovations = [
    ('动态动画生成引擎',
     '不同于固定的视频或 GIF，系统后端 AnimationService 根据算法类型动态\n生成每步的动画数据（比较索引、交换位置、状态变化等）。',
     '确保动画数据与算法定义一致，新增算法只需配置伪代码即可自动生成动画。'),
    ('伪代码同步高亮',
     '动画每推进一步，对应的伪代码行同步高亮。\n学生可以清楚地看到"这一步在执行哪行代码"。',
     '将抽象代码与具体操作建立直观的因果关系。'),
    ('全角色权限体系',
     '从零构建了三层角色（student/teacher/admin）的\n完整权限体系，包含 15 个管理员专属 API。',
     '管理员可在网页端直接管理用户、算法和学习内容，无需数据库操作。'),
    ('工程化全栈实践',
     'React 18 + TypeScript + Vite + TailwindCSS + Zustand\nPython Flask + PostgreSQL + JWT + SQLAlchemy',
     '前后端分离、类型安全、状态持久化、暗色主题、响应式布局 — 接近生产级代码质量。'),
]
for i, (title, desc, impact) in enumerate(innovations):
    y = 1.6 + i * 1.4
    add_shape_rect(slide, 1.2, y, 11, 1.2, fill_color=DARK_CARD)
    add_text_box(slide, 1.5, y + 0.1, 3, 0.35, title, font_size=18, color=PRIMARY, bold=True)
    add_text_box(slide, 1.5, y + 0.5, 10, 0.6, desc, font_size=13, color=WHITE)
    add_text_box(slide, 1.5, y + 0.85, 10, 0.3, f'💡 {impact}', font_size=12, color=LIGHT_GRAY)

# ============================================================
# SLIDE 19 — 不足与展望
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, '不足与展望', '3')

add_text_box(slide, 1.2, 1.6, 5, 0.5, '存在的不足', font_size=22, color=RED, bold=True)
weaknesses = [
    '• 动画类型覆盖有限，目前支持 5 大类 9 种算法',
    '• 缺少单元测试和集成测试',
    '• 未实现 Redis 缓存（依赖已安装但未启用）',
    '• 缺少国际化支持，目前仅支持中文',
    '• 管理后台未实现算法动画数据的可视化编辑',
]
for i, w in enumerate(weaknesses):
    add_text_box(slide, 1.2, 2.3 + i * 0.5, 5.5, 0.45, w, font_size=14, color=LIGHT_GRAY)

add_text_box(slide, 7.2, 1.6, 5, 0.5, '未来展望', font_size=22, color=GREEN, bold=True)
futures = [
    '• 扩展更多算法类型（贪心、回溯、字符串匹配等）',
    '• 增加在线代码编辑与执行功能',
    '• 引入 WebSocket 实现协作学习功能',
    '• 集成 AI 辅助学习（LLM 答疑）',
    '• 部署上线，支持移动端适配',
    '• 增加数据可视化分析（学习行为画像）',
]
for i, f in enumerate(futures):
    add_text_box(slide, 7.2, 2.3 + i * 0.5, 5.5, 0.45, f, font_size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 20 — Section: 致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
# Decorative top bar
dec = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
dec.fill.solid(); dec.fill.fore_color.rgb = PRIMARY; dec.line.fill.background()

add_text_box(slide, 1.5, 2.0, 10, 1.0, '致  谢', font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_bar(slide, 5.8, 3.0, 1.5)

thanks_text = (
    '感谢指导教师在毕业设计过程中给予的悉心指导与宝贵建议。\n\n'
    '感谢各位评审老师在百忙之中审阅本论文并参加答辩。\n\n'
    '感谢同学和朋友们的帮助与支持。\n\n'
    '感谢所有开源项目与社区贡献者。'
)
add_text_box(slide, 2.5, 3.5, 8, 3.0, thanks_text, font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 2.0, 6.5, 9, 0.5, '请各位老师批评指正！', font_size=20, color=PRIMARY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = r'C:\Users\啊\Desktop\bi_ye_she_ji_2\毕业设计答辩PPT.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
